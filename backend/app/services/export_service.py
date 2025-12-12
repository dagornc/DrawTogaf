from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from io import BytesIO
from typing import Dict, List, Any
import logging

logger = logging.getLogger(__name__)

class ExportService:
    def __init__(self):
        pass

    def _get_color_by_type(self, type_str: str):
        t = type_str.lower() if type_str else ""
        
        # Colors (RGB) matching the frontend
        # Active (Blue): bg #E1F5FE (225, 245, 254), border #0288D1 (2, 136, 209)
        # Behavior (Yellow): bg #FFF9C4 (255, 249, 196), border #FBC02D (251, 192, 45)
        # Passive (Green): bg #E8F5E9 (232, 245, 233), border #43A047 (67, 160, 71)

        if t in [
            'applicationcomponent', 'applicationinterface',
            'businessactor', 'businessrole',
            'node', 'device', 'systemsoftware', 'communicationnetwork'
        ]:
            return RGBColor(225, 245, 254), RGBColor(2, 136, 209) # Active
        
        if t in [
            'businessprocess', 'businessfunction', 'businessinteraction', 'businessservice',
            'applicationservice', 'applicationfunction', 'applicationinteraction',
            'technologyservice', 'technologyfunction'
        ]:
            return RGBColor(255, 249, 196), RGBColor(251, 192, 45) # Behavior
            
        if t in [
            'dataobject', 'businessobject', 'artifact', 'contract', 'representation'
        ]:
            return RGBColor(232, 245, 233), RGBColor(67, 160, 71) # Passive

        return RGBColor(245, 245, 245), RGBColor(158, 158, 158) # Default

    def create_pptx(self, graph_data: Dict[str, Any]) -> BytesIO:
        """
        Generates a PowerPoint file from the graph data.
        """
        prs = Presentation()
        # Use a blank slide layout (usually index 6)
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)
        
        nodes = graph_data.get("nodes", [])
        edges = graph_data.get("edges", [])
        
        # Determine diagram bounds
        if not nodes:
            # Empty diagram, valid pptx but empty
            output = BytesIO()
            prs.save(output)
            output.seek(0)
            return output

        # Extract coordinates and dimensions, falling back to 0 if missing (though they shouldn't be now)
        # We need to handle potential string values or missing keys safely
        xs = [float(n.get("position", {}).get("x", 0) or 0) for n in nodes]
        ys = [float(n.get("position", {}).get("y", 0) or 0) for n in nodes]
        
        # Calculate bounding box of all nodes including their width/height
        # to ensure nothing is clipped
        min_x = min(xs)
        min_y = min(ys)
        
        # Calculate max extents
        max_x = max([float(n.get("position", {}).get("x", 0) or 0) + float(n.get("width", 150) or 150) for n in nodes])
        max_y = max([float(n.get("position", {}).get("y", 0) or 0) + float(n.get("height", 80) or 80) for n in nodes])
        
        diagram_width = max_x - min_x
        diagram_height = max_y - min_y
        
        # Protect against div by zero for single point diagrams
        if diagram_width == 0: diagram_width = 1
        if diagram_height == 0: diagram_height = 1

        # PowerPoint Slide Dimensions (Standard 4:3 is 10x7.5 inches, Widescreen 16:9 is 13.33x7.5 inches)
        # python-pptx default is 10x7.5 inches (9144000 EMUs width, 6858000 EMUs height)
        # 1 inch = 914400 EMUs. 
        # Slide width in points: 10 * 72 = 720 pt
        # Slide height in points: 7.5 * 72 = 540 pt
        
        SLIDE_WIDTH_PT = prs.slide_width.pt
        SLIDE_HEIGHT_PT = prs.slide_height.pt
        
        # Margins (50pt)
        MARGIN = 50
        available_width = SLIDE_WIDTH_PT - (2 * MARGIN)
        available_height = SLIDE_HEIGHT_PT - (2 * MARGIN)
        
        # Calculate Scale Factor to fit diagram into available space
        scale_x = available_width / diagram_width
        scale_y = available_height / diagram_height
        scale = min(scale_x, scale_y)
        
        # Don't scale up if the diagram is small, just center it? 
        # Actually scaling up is fine for visibility, but maybe cap it at 1.0 or 1.5 to check?
        # User said "see all objects", so scaling down is the priority. Scaling up is a bonus.
        # Let's simple use the calculated scale.
        
        # Center the diagram
        # The scaled diagram dimensions
        scaled_w = diagram_width * scale
        scaled_h = diagram_height * scale
        
        # Offsets to center
        offset_x = MARGIN + (available_width - scaled_w) / 2
        offset_y = MARGIN + (available_height - scaled_h) / 2
        
        node_shapes = {}
        
        for node in nodes:
            # Original coords
            x = float(node.get("position", {}).get("x", 0) or 0)
            y = float(node.get("position", {}).get("y", 0) or 0)
            w = float(node.get("width", 150) or 150)
            h = float(node.get("height", 80) or 80)
            
            label = node.get("name") or node.get("data", {}).get("label") or node.get("label", "Node")
            node_type = node.get("type") or node.get("data", {}).get("type") or "Unknown"
            
            # Use layer to override type for color if specific generic type is used
            # But the existing color logic relies on type strings.
            
            # Apply scaling
            # New coordinate = (Original - Min) * Scale + Offset
            left_pt = (x - min_x) * scale + offset_x
            top_pt = (y - min_y) * scale + offset_y
            width_pt = w * scale
            height_pt = h * scale
            
            fill_color, line_color = self._get_color_by_type(node_type)
            
            # Create Shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt)
            )
            
            # Style
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.line.color.rgb = line_color
            shape.line.width = Pt(1.5) # This stays constant width regardless of scale? Or logic says thicker? constant is fine.
            
            # Text
            text_frame = shape.text_frame
            text_frame.text = label
            # Adjust font size based on scale to avoid huge text on small blocks
            # Base font 12?
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(max(8, 12 * scale)) # Minimum 8pt
                    
            node_shapes[node["id"]] = shape
            
        # Draw Edges
        for edge in edges:
            source_id = edge.get("source_id") or edge.get("source")
            target_id = edge.get("target_id") or edge.get("target")
            
            if source_id in node_shapes and target_id in node_shapes:
                source_shape = node_shapes[source_id]
                target_shape = node_shapes[target_id]
                
                connector = slide.shapes.add_connector(
                    MSO_SHAPE.STRAIGHT_CONNECTOR_1, 0, 0, 0, 0
                )
                
                connector.begin_connect(source_shape, idx=2) # Bottom
                connector.end_connect(target_shape, idx=0)   # Top
                
                connector.line.color.rgb = RGBColor(100, 100, 100)
                connector.line.width = Pt(1)
                
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output
